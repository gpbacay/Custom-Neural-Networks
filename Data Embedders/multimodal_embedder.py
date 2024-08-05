import tensorflow as tf
from transformers import TFBertModel, BertTokenizer
from tensorflow.keras.applications import ResNet50, resnet50
from tensorflow.keras.preprocessing import image as keras_image
from tensorflow.keras.applications.resnet50 import preprocess_input
from tensorflow.keras.datasets import mnist
import numpy as np
from PIL import Image
import requests
from io import BytesIO
import librosa
from moviepy.editor import VideoFileClip
from docx import Document
import PyPDF2
from pptx import Presentation
import pandas as pd

class DataEmbedder:
    def __init__(self):
        # Initialize text model and tokenizer
        self.text_model = TFBertModel.from_pretrained('bert-base-uncased')
        self.text_tokenizer = BertTokenizer.from_pretrained('bert-base-uncased')

        # Initialize image model
        self.image_model = ResNet50(weights='imagenet', include_top=False, pooling='avg')

    def embed_text(self, text):
        inputs = self.text_tokenizer(text, return_tensors='tf', truncation=True, padding=True)
        outputs = self.text_model(**inputs)
        embeddings = tf.reduce_mean(outputs.last_hidden_state, axis=1)
        return embeddings

    def embed_image(self, image_path):
        img = keras_image.load_img(image_path, target_size=(224, 224))
        img_array = keras_image.img_to_array(img)
        img_array = np.expand_dims(img_array, axis=0)
        img_array = preprocess_input(img_array)
        embeddings = self.image_model.predict(img_array)
        return embeddings

    def embed_audio(self, audio_path):
        y, sr = librosa.load(audio_path, sr=None)
        mel_spectrogram = librosa.feature.melspectrogram(y=y, sr=sr)
        log_mel_spectrogram = librosa.power_to_db(mel_spectrogram, ref=np.max)
        log_mel_spectrogram = np.expand_dims(log_mel_spectrogram, axis=-1)
        log_mel_spectrogram = np.expand_dims(log_mel_spectrogram, axis=0)
        log_mel_spectrogram = tf.image.resize(log_mel_spectrogram, [224, 224])
        log_mel_spectrogram = np.array(log_mel_spectrogram)
        embeddings = self.image_model.predict(log_mel_spectrogram)
        return embeddings

    def embed_video(self, video_path):
        clip = VideoFileClip(video_path)
        # Sample frames from the video
        frames = []
        for frame in clip.iter_frames(fps=1):  # Sample one frame per second
            img = Image.fromarray(frame)
            img = img.resize((224, 224))
            img_array = keras_image.img_to_array(img)
            img_array = preprocess_input(img_array)
            frames.append(img_array)
        frames = np.array(frames)
        embeddings = self.image_model.predict(frames)
        return np.mean(embeddings, axis=0)  # Average the frame embeddings

    def embed_mnist(self, mnist_image):
        mnist_image = mnist_image.reshape((28, 28))  # Ensure the image is 28x28
        mnist_image = np.stack([mnist_image] * 3, axis=-1)  # Convert to 3 channels
        mnist_image = Image.fromarray(mnist_image.astype('uint8')).resize((224, 224))
        img_array = keras_image.img_to_array(mnist_image)
        img_array = np.expand_dims(img_array, axis=0)
        img_array = preprocess_input(img_array)
        embeddings = self.image_model.predict(img_array)
        return embeddings

    def read_text_file(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read()

    def read_docx_file(self, file_path):
        doc = Document(file_path)
        return '\n'.join([para.text for para in doc.paragraphs])

    def read_pdf_file(self, file_path):
        pdf_reader = PyPDF2.PdfReader(file_path)
        text = ''
        for page in pdf_reader.pages:
            text += page.extract_text()
        return text

    def read_pptx_file(self, file_path):
        prs = Presentation(file_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        return text

    def read_csv_file(self, file_path):
        df = pd.read_csv(file_path)
        return df.to_string()

    def read_excel_file(self, file_path):
        df = pd.read_excel(file_path)
        return df.to_string()

    def detect_and_embed(self, input_data):
        # Detect data type
        if isinstance(input_data, str):
            if input_data.endswith(('.png', '.jpg', '.jpeg')):
                # It's an image URL or path
                return self.embed_image(input_data)
            elif input_data.endswith(('.wav', '.mp3')):
                # It's an audio file path
                return self.embed_audio(input_data)
            elif input_data.endswith('.mp4'):
                # It's a video file path
                return self.embed_video(input_data)
            elif input_data.endswith('.txt'):
                # It's a text file
                text = self.read_text_file(input_data)
                return self.embed_text(text)
            elif input_data.endswith('.docx'):
                # It's a DOCX file
                text = self.read_docx_file(input_data)
                return self.embed_text(text)
            elif input_data.endswith('.pdf'):
                # It's a PDF file
                text = self.read_pdf_file(input_data)
                return self.embed_text(text)
            elif input_data.endswith('.pptx'):
                # It's a PPTX file
                text = self.read_pptx_file(input_data)
                return self.embed_text(text)
            elif input_data.endswith('.csv'):
                # It's a CSV file
                text = self.read_csv_file(input_data)
                return self.embed_text(text)
            elif input_data.endswith('.xlsx'):
                # It's an Excel file
                text = self.read_excel_file(input_data)
                return self.embed_text(text)
            else:
                # Assume it's plain text
                return self.embed_text(input_data)
        elif isinstance(input_data, bytes):
            # Assume it's image data
            img = Image.open(BytesIO(input_data))
            img.save("temp_image.jpg")
            return self.embed_image("temp_image.jpg")
        elif isinstance(input_data, np.ndarray) and input_data.shape == (28, 28):
            # Assume it's an MNIST image
            return self.embed_mnist(input_data)
        else:
            raise ValueError("Unsupported data type")

# Example usage:
embedder = DataEmbedder()

# # Embedding text
# text = "This is an example sentence."
# text_embedding = embedder.detect_and_embed(text)
# print("Text Embedding:", text_embedding)

# # Embedding image
# image_url = "https://example.com/image.jpg"
# response = requests.get(image_url)
# image_embedding = embedder.detect_and_embed(response.content)
# print("Image Embedding:", image_embedding)

# # Embedding audio
# audio_path = "path_to_audio_file.mp3"  # Replace with your audio file path
# audio_embedding = embedder.detect_and_embed(audio_path)
# print("Audio Embedding:", audio_embedding)

# # Embedding video
# video_path = "path_to_video_file.mp4"  # Replace with your video file path
# video_embedding = embedder.detect_and_embed(video_path)
# print("Video Embedding:", video_embedding)

# Embedding MNIST data
(mnist_train_images, mnist_train_labels), (mnist_test_images, mnist_test_labels) = mnist.load_data()
mnist_image = mnist_test_images[0]
mnist_embedding = embedder.detect_and_embed(mnist_image)
print("MNIST Embedding:", mnist_embedding)

# # Embedding .txt file
# txt_file_path = "path_to_text_file.txt"  # Replace with your text file path
# txt_embedding = embedder.detect_and_embed(txt_file_path)
# print("Text File Embedding:", txt_embedding)

# # Embedding .docx file
# docx_file_path = "path_to_docx_file.docx"  # Replace with your DOCX file path
# docx_embedding = embedder.detect_and_embed(docx_file_path)
# print("DOCX File Embedding:", docx_embedding)

# # Embedding .pdf file
# pdf_file_path = "path_to_pdf_file.pdf"  # Replace with your PDF file path
# pdf_embedding = embedder.detect_and_embed(pdf_file_path)
# print("PDF File Embedding:", pdf_embedding)

# # Embedding .pptx file
# pptx_file_path = "path_to_pptx_file.pptx"  # Replace with your PPTX file path
# pptx_embedding = embedder.detect_and_embed(pptx_file_path)
# print("PPTX File Embedding:", pptx_embedding)

# # Embedding .csv file
# csv_file_path = "path_to_csv_file.csv"  # Replace with your CSV file path
# csv_embedding = embedder.detect_and_embed(csv_file_path)
# print("CSV File Embedding:", csv_embedding)

# # Embedding .xlsx file
# xlsx_file_path = "path_to_excel_file.xlsx"  # Replace with your Excel file path
# xlsx_embedding = embedder.detect_and_embed(xlsx_file_path)
# print("Excel File Embedding:", xlsx_embedding)
